"""
Document Processing Pipeline
Handles file uploads, text extraction, OCR, and document classification
"""

import os
import io
import base64
from typing import List, Dict, Tuple, Optional
from pathlib import Path
import logging

import PyPDF2
import pdfplumber
from docx import Document
from PIL import Image
import pandas as pd
from pptx import Presentation
import requests
from tqdm import tqdm

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Main document processing class"""
    
    SUPPORTED_FORMATS = ['.pdf', '.docx', '.txt', '.csv', '.xlsx', '.pptx', '.jpg', '.jpeg', '.png']
    
    def __init__(self, hf_api_key: str, ocr_model: str, classification_model: str):
        self.hf_api_key = hf_api_key
        self.ocr_model = ocr_model
        self.classification_model = classification_model
        self.hf_api_headers = {"Authorization": f"Bearer {hf_api_key}"}
        
    def process_document(self, file_path: str) -> Dict:
        """
        Main entry point for document processing
        
        Returns:
            Dict with extracted_text, metadata, document_type, pages
        """
        try:
            file_path = Path(file_path)
            extension = file_path.suffix.lower()
            
            if extension not in self.SUPPORTED_FORMATS:
                raise ValueError(f"Unsupported file format: {extension}")
            
            logger.info(f"Processing document: {file_path.name}")
            
            # Extract text based on file type
            if extension == '.pdf':
                extracted_data = self._process_pdf(file_path)
            elif extension == '.docx':
                extracted_data = self._process_docx(file_path)
            elif extension == '.txt':
                extracted_data = self._process_txt(file_path)
            elif extension in ['.csv', '.xlsx']:
                extracted_data = self._process_spreadsheet(file_path)
            elif extension == '.pptx':
                extracted_data = self._process_pptx(file_path)
            elif extension in ['.jpg', '.jpeg', '.png']:
                extracted_data = self._process_image(file_path)
            else:
                raise ValueError(f"Handler not implemented for {extension}")
            
            # Classify document type
            document_type = self._classify_document(extracted_data['extracted_text'])
            extracted_data['document_type'] = document_type
            
            logger.info(f"Document classified as: {document_type}")
            
            return extracted_data
            
        except Exception as e:
            logger.error(f"Error processing document: {e}")
            raise
    
    def _process_pdf(self, file_path: Path) -> Dict:
        """Extract text from PDF, with OCR fallback for image-based pages"""
        extracted_text = ""
        pages = []
        metadata = {
            'filename': file_path.name,
            'file_type': 'pdf',
            'total_pages': 0
        }
        
        try:
            # First attempt: Try extracting text directly
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                metadata['total_pages'] = len(pdf_reader.pages)
                
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    page_text = page.extract_text() or ""
                    
                    # If very little text extracted, might be image-based
                    if len(page_text.strip()) < 50:
                        logger.info(f"Page {page_num} appears to be image-based, attempting OCR...")
                        # Use pdfplumber to extract images and perform OCR
                        page_text = self._ocr_pdf_page(file_path, page_num - 1)
                    
                    pages.append({
                        'page_number': page_num,
                        'text': page_text
                    })
                    extracted_text += f"\n[Page {page_num}]\n{page_text}\n"
            
        except Exception as e:
            logger.error(f"Error processing PDF: {e}")
            # Fallback to pdfplumber
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page_num, page in enumerate(pdf.pages, 1):
                        page_text = page.extract_text() or ""
                        pages.append({
                            'page_number': page_num,
                            'text': page_text
                        })
                        extracted_text += f"\n[Page {page_num}]\n{page_text}\n"
            except Exception as e2:
                logger.error(f"Fallback PDF processing failed: {e2}")
                raise
        
        return {
            'extracted_text': extracted_text.strip(),
            'metadata': metadata,
            'pages': pages
        }
    
    def _ocr_pdf_page(self, file_path: Path, page_num: int) -> str:
        """Perform OCR on a specific PDF page using Hugging Face API"""
        try:
            # Convert PDF page to image
            with pdfplumber.open(file_path) as pdf:
                if page_num < len(pdf.pages):
                    page = pdf.pages[page_num]
                    # Convert to image
                    img = page.to_image(resolution=300)
                    img_pil = img.original
                    
                    # Perform OCR
                    ocr_text = self._perform_ocr(img_pil)
                    return ocr_text
        except Exception as e:
            logger.error(f"OCR failed for page {page_num}: {e}")
        
        return ""
    
    def _perform_ocr(self, image: Image.Image) -> str:
        """Perform OCR using Hugging Face TrOCR API"""
        try:
            # Convert PIL Image to bytes
            img_byte_arr = io.BytesIO()
            image.save(img_byte_arr, format='PNG')
            img_byte_arr = img_byte_arr.getvalue()
            
            # Call Hugging Face API
            api_url = f"https://api-inference.huggingface.co/models/{self.ocr_model}"
            response = requests.post(
                api_url,
                headers=self.hf_api_headers,
                data=img_byte_arr,
                timeout=30
            )
            
            if response.status_code == 200:
                result = response.json()
                # TrOCR returns a list with generated_text
                if isinstance(result, list) and len(result) > 0:
                    return result[0].get('generated_text', '')
                elif isinstance(result, dict):
                    return result.get('generated_text', '')
            else:
                logger.warning(f"OCR API returned status {response.status_code}")
                
        except Exception as e:
            logger.error(f"OCR API call failed: {e}")
        
        return ""
    
    def _process_docx(self, file_path: Path) -> Dict:
        """Extract text from DOCX files"""
        extracted_text = ""
        pages = []
        
        try:
            doc = Document(file_path)
            
            for idx, para in enumerate(doc.paragraphs):
                if para.text.strip():
                    extracted_text += para.text + "\n"
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        extracted_text += cell.text + " "
                extracted_text += "\n"
            
            pages.append({'page_number': 1, 'text': extracted_text})
            
        except Exception as e:
            logger.error(f"Error processing DOCX: {e}")
            raise
        
        metadata = {
            'filename': file_path.name,
            'file_type': 'docx',
            'total_pages': 1
        }
        
        return {
            'extracted_text': extracted_text.strip(),
            'metadata': metadata,
            'pages': pages
        }
    
    def _process_txt(self, file_path: Path) -> Dict:
        """Extract text from TXT files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                extracted_text = file.read()
        except UnicodeDecodeError:
            # Try different encoding
            with open(file_path, 'r', encoding='latin-1') as file:
                extracted_text = file.read()
        
        metadata = {
            'filename': file_path.name,
            'file_type': 'txt',
            'total_pages': 1
        }
        
        return {
            'extracted_text': extracted_text.strip(),
            'metadata': metadata,
            'pages': [{'page_number': 1, 'text': extracted_text}]
        }
    
    def _process_spreadsheet(self, file_path: Path) -> Dict:
        """Extract text from CSV/Excel files"""
        try:
            if file_path.suffix.lower() == '.csv':
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path)
            
            # Convert to text representation
            extracted_text = df.to_string()
            
            metadata = {
                'filename': file_path.name,
                'file_type': file_path.suffix.lower().replace('.', ''),
                'total_pages': 1,
                'rows': len(df),
                'columns': len(df.columns)
            }
            
            return {
                'extracted_text': extracted_text.strip(),
                'metadata': metadata,
                'pages': [{'page_number': 1, 'text': extracted_text}]
            }
        except Exception as e:
            logger.error(f"Error processing spreadsheet: {e}")
            raise
    
    def _process_pptx(self, file_path: Path) -> Dict:
        """Extract text from PowerPoint files"""
        extracted_text = ""
        pages = []
        
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                
                pages.append({
                    'page_number': slide_num,
                    'text': slide_text
                })
                extracted_text += f"\n[Slide {slide_num}]\n{slide_text}\n"
            
            metadata = {
                'filename': file_path.name,
                'file_type': 'pptx',
                'total_pages': len(prs.slides)
            }
            
            return {
                'extracted_text': extracted_text.strip(),
                'metadata': metadata,
                'pages': pages
            }
        except Exception as e:
            logger.error(f"Error processing PPTX: {e}")
            raise
    
    def _process_image(self, file_path: Path) -> Dict:
        """Extract text from images using OCR"""
        try:
            image = Image.open(file_path)
            
            # Perform OCR
            extracted_text = self._perform_ocr(image)
            
            metadata = {
                'filename': file_path.name,
                'file_type': 'image',
                'total_pages': 1,
                'image_size': image.size
            }
            
            return {
                'extracted_text': extracted_text.strip(),
                'metadata': metadata,
                'pages': [{'page_number': 1, 'text': extracted_text}]
            }
        except Exception as e:
            logger.error(f"Error processing image: {e}")
            raise
    
    def _classify_document(self, text: str) -> str:
        """
        Classify document type using keyword detection + zero-shot classification
        Categories: Job Description, Resume, Research Paper, Legal Document, Invoice, Textbook, Generic
        """
        if not text or len(text.strip()) < 50:
            return "Generic Document"
        
        # First try keyword-based detection for common types
        text_lower = text.lower()
        
        # Job Description keywords (more weighted and specific)
        jd_keywords = [
            'job description', 'responsibilities', 'requirements', 'qualifications',
            'job role', 'position', 'vacancy', 'hiring', 'candidate', 'apply',
            'years of experience', 'skills required', 'we are looking for',
            'job summary', 'duties', 'preferred qualifications', 'employment type',
            'key responsibilities', 'seeking a', 'join our team', 'about the role',
            'what you will do', 'required experience', 'desired skills',
            'compensation', 'benefits package', 'equal opportunity employer'
        ]
        jd_score = sum(2 if keyword in text_lower else 0 for keyword in jd_keywords)
        
        # Resume keywords (personal pronouns are strong indicators)
        resume_keywords = [
            'education', 'projects', 'achievements', 'my experience', 'i have',
            'certifications', 'internship', 'objective', 'profile summary',
            'work experience', 'bachelor', 'master', 'degree', 'university',
            'technical skills', 'linkedin', 'github', 'portfolio', 'contact',
            'email', 'phone', 'address', 'references available'
        ]
        resume_score = sum(1 for keyword in resume_keywords if keyword in text_lower)
        
        # Boost JD score if it mentions "seeking" or "looking for" without personal pronouns
        if any(phrase in text_lower for phrase in ['seeking a', 'looking for a', 'we are seeking', 'we are looking']):
            jd_score += 5
        
        # Reduce Resume score if it has company hiring language
        if any(phrase in text_lower for phrase in ['we are seeking', 'join our team', 'about the role']):
            resume_score = max(0, resume_score - 5)
        
        # Research Paper keywords
        research_keywords = [
            'abstract', 'introduction', 'methodology', 'results', 'conclusion',
            'references', 'citation', 'journal', 'published', 'doi', 'arxiv',
            'literature review', 'hypothesis', 'experiment', 'analysis'
        ]
        research_score = sum(1 for keyword in research_keywords if keyword in text_lower)
        
        # Legal Document keywords
        legal_keywords = [
            'agreement', 'contract', 'whereas', 'hereby', 'party', 'clause',
            'terms and conditions', 'liability', 'jurisdiction', 'legal',
            'attorney', 'law', 'court', 'plaintiff', 'defendant'
        ]
        legal_score = sum(1 for keyword in legal_keywords if keyword in text_lower)
        
        # Invoice keywords
        invoice_keywords = [
            'invoice', 'bill', 'amount', 'payment', 'total', 'tax',
            'subtotal', 'due date', 'invoice number', 'billing address',
            'quantity', 'price', 'gst', 'discount'
        ]
        invoice_score = sum(1 for keyword in invoice_keywords if keyword in text_lower)
        
        # Determine highest scoring category
        scores = {
            'Job Description': jd_score,
            'Resume or CV': resume_score,
            'Research Paper or Academic Article': research_score,
            'Legal Document or Contract': legal_score,
            'Invoice or Financial Report': invoice_score
        }
        
        # If any category has strong signals (>= 5 keywords), use it
        max_category = max(scores, key=scores.get)
        if scores[max_category] >= 5:
            logger.info(f"Document classified via keywords as: {max_category} (score: {scores[max_category]})")
            return max_category
        
        # Otherwise, fall back to API-based classification
        candidate_labels = [
            "Job Description",
            "Resume or CV",
            "Research Paper or Academic Article",
            "Legal Document or Contract",
            "Invoice or Financial Report",
            "Textbook or Educational Material",
            "Generic Document"
        ]
        
        try:
            # Truncate text for API call (max ~500 words)
            text_sample = " ".join(text.split()[:500])
            
            api_url = f"https://api-inference.huggingface.co/models/{self.classification_model}"
            payload = {
                "inputs": text_sample,
                "parameters": {
                    "candidate_labels": candidate_labels
                }
            }
            
            response = requests.post(
                api_url,
                headers=self.hf_api_headers,
                json=payload,
                timeout=30
            )
            
            if response.status_code == 200:
                result = response.json()
                # Get the label with highest score
                if 'labels' in result and len(result['labels']) > 0:
                    return result['labels'][0]
            else:
                logger.warning(f"Classification API returned status {response.status_code}")
                
        except Exception as e:
            logger.error(f"Document classification failed: {e}")
        
        return "Generic Document"


class DocumentClassifier:
    """Dedicated document classifier for better organization"""
    
    DOCUMENT_TYPES = {
        "Job Description": "job_description",
        "Resume or CV": "resume",
        "Research Paper or Academic Article": "research_paper",
        "Legal Document or Contract": "legal",
        "Invoice or Financial Report": "invoice",
        "Textbook or Educational Material": "textbook",
        "Generic Document": "generic"
    }
    
    @staticmethod
    def get_document_category(doc_type: str) -> str:
        """Convert full document type to category"""
        return DocumentClassifier.DOCUMENT_TYPES.get(doc_type, "generic")
